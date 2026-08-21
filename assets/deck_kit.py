# -*- coding: utf-8 -*-
"""生物医学精读幻灯片 · 设计系统工具库
米白底 + 深海军蓝 + 暗砖红 | 微软雅黑 + Arial | 趋势对照式表达
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
from PIL import Image

NAVY      = RGBColor(0x1F, 0x38, 0x64)
NAVY_DEEP = RGBColor(0x16, 0x28, 0x48)
RED       = RGBColor(0xA6, 0x3D, 0x2F)
BG        = RGBColor(0xFB, 0xFA, 0xF6)
INK       = RGBColor(0x2B, 0x2B, 0x2B)
GRAY      = RGBColor(0x6E, 0x6E, 0x6E)
LIGHT_RED = RGBColor(0xF7, 0xE9, 0xE4)
LIGHT_BLU = RGBColor(0xEA, 0xEF, 0xF5)
BORDER    = RGBColor(0xC9, 0xCF, 0xD6)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
SOFT_WHT  = RGBColor(0xC9, 0xD4, 0xE4)

FONT_EA = "微软雅黑"
FONT_LATIN = "Arial"

STYLES = {
    "n":  dict(color=INK,  bold=False),
    "b":  dict(color=INK,  bold=True),
    "r":  dict(color=RED,  bold=False),
    "rb": dict(color=RED,  bold=True),
    "bl": dict(color=NAVY, bold=True),
    "g":  dict(color=GRAY, bold=False),
}

EMU_IN = 914400
SLIDE_W = 13.333
SLIDE_H = 7.5

def make_prs():
    prs = Presentation()
    prs.slide_width = Emu(int(SLIDE_W * EMU_IN))
    prs.slide_height = Emu(int(SLIDE_H * EMU_IN))
    return prs

def new_slide(prs, bg=BG):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s

def _style_run(run, size, color, bold, italic=False):
    f = run.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color; f.name = FONT_LATIN
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", FONT_EA)

def add_text(slide, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p.get("align", align)
        if p.get("space_after") is not None: para.space_after = Pt(p["space_after"])
        if p.get("space_before") is not None: para.space_before = Pt(p["space_before"])
        if p.get("line") is not None: para.line_spacing = p["line"]
        base_size = p.get("size", 12)
        for item in p["runs"]:
            text, style = item[0], item[1]
            sz = item[2] if len(item) > 2 else base_size
            run = para.add_run(); run.text = text
            st = STYLES[style]
            _style_run(run, sz, p.get("color", st["color"]),
                       st["bold"] or p.get("bold", False), p.get("italic", False))
            if p.get("color") is not None:
                run.font.color.rgb = p["color"]
    return tb

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75,
             shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None: sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0] = radius
        except Exception: pass
    return sp

def add_pic_fit(slide, path, x, y, max_w, max_h, border=True):
    iw, ih = Image.open(path).size
    scale = min(max_w / iw, max_h / ih)
    w, h = iw * scale, ih * scale
    px = x + (max_w - w) / 2
    pic = slide.shapes.add_picture(path, Inches(px), Inches(y), Inches(w), Inches(h))
    if border:
        pic.line.color.rgb = BORDER; pic.line.width = Pt(1.0)
    return pic, h

def page_number(slide, idx, dark_bg=False):
    add_text(slide, 12.75, 7.14, 0.45, 0.3,
             [dict(runs=[(str(idx), "n")], size=9,
                   color=(SOFT_WHT if dark_bg else GRAY))], align=PP_ALIGN.RIGHT)

def title_bar(slide, title_runs, chip_text=None):
    add_rect(slide, 0, 0, SLIDE_W, 1.0, fill=NAVY)
    add_rect(slide, 0, 1.0, SLIDE_W, 0.035, fill=RED)
    add_text(slide, 0.45, 0, 9.6, 1.0,
             [dict(runs=title_runs, size=20, color=WHITE, bold=True)],
             anchor=MSO_ANCHOR.MIDDLE)
    if chip_text:
        chip_w = 2.75
        add_rect(slide, SLIDE_W - 0.45 - chip_w, 0.30, chip_w, 0.42,
                 fill=RED, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
        add_text(slide, SLIDE_W - 0.45 - chip_w, 0.30, chip_w, 0.42,
                 [dict(runs=[(chip_text, "n")], size=11.5, color=WHITE, bold=True)],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def takeaway_bar(slide, text_runs, label="结果意义"):
    add_rect(slide, 0, 6.78, SLIDE_W, 0.56, fill=LIGHT_RED)
    add_rect(slide, 0, 6.78, 0.12, 0.56, fill=RED)
    runs = [(label + "  ", "rb", 12.5)] + [(t, s, 12.5) for t, s in text_runs]
    add_text(slide, 0.45, 6.78, 12.5, 0.56, [dict(runs=runs, size=12.5)],
             anchor=MSO_ANCHOR.MIDDLE)

def section_paras(sections):
    paras = []
    for head, bullets in sections:
        paras.append(dict(runs=[("▍ ", "rb", 13.5), (head, "bl", 13.5)],
                          size=13.5, space_before=13, space_after=5))
        for b in bullets:
            runs = [("•  ", "r", 12.5)] + [(t, s, 12.5) for t, s in b]
            paras.append(dict(runs=runs, size=12.5, space_after=5, line=1.22))
    return paras

def _mini_panel(slide, x, y, w, h, label, bullets):
    """顶部信息小面板：为什么做 / 做了什么"""
    add_rect(slide, x, y, w, h, fill=LIGHT_BLU,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
    paras = [dict(runs=[(label, "n")], size=11, color=NAVY, bold=True, space_after=3)]
    for b in bullets:
        runs = [("•  ", "r", 10.5)] + [(t, st, 10.5) for t, st in b]
        paras.append(dict(runs=runs, size=10.5, space_after=2, line=1.12))
    add_text(slide, x + 0.22, y + 0.11, w - 0.44, h - 0.22, paras)

def _place_image(slide, path, region_x, region_y, region_w, region_h,
                 caption, cap_align=PP_ALIGN.CENTER):
    """在指定区域内尽可能大地放置图片(含图注)，整体垂直水平居中"""
    iw, ih = Image.open(path).size
    box_h = region_h - 0.32  # 给图注留位
    w = min(region_w, box_h * iw / ih)
    h = w * ih / iw
    block = h + 0.30
    y0 = region_y + max(0.0, (region_h - block) / 2)
    x0 = region_x + max(0.0, (region_w - w) / 2)
    pic = slide.shapes.add_picture(path, Inches(x0), Inches(y0), Inches(w), Inches(h))
    pic.line.color.rgb = BORDER
    pic.line.width = Pt(1.0)
    add_text(slide, region_x, y0 + h + 0.06, region_w, 0.3,
             [dict(runs=[(caption, "g")], size=9.5, align=cap_align)])

def _results_head(size=12.0):
    return dict(runs=[("▍ ", "rb", size), ("结果分析（趋势对比）", "bl", size)],
                size=size, space_after=5)

def _result_bullet(b, size=11.0):
    runs = [("•  ", "r", size)] + [(t, st, size) for t, st in b]
    return dict(runs=runs, size=size, space_after=4, line=1.2)

def experiment_slide(prs, img_dir, idx, chip, title, fig,
                     why, how, results, takeaway, img, caption):
    """以图为中心的自适应布局：
    宽幅图(r>=2.2): 顶部 why/how → 全宽大图 → 结果在图下
    中高图(r<2.2) : 顶部 why/how → 左侧满高大图 → 结果在图右
    """
    s = new_slide(prs)
    title_bar(s, [(title, "n", 20), ("   " + fig, "n", 15)], chip_text=chip)
    _mini_panel(s, 0.45, 1.12, 6.10, 1.0, "为什么做", why)
    _mini_panel(s, 6.80, 1.12, 6.10, 1.0, "做了什么", how)
    path = os.path.join(img_dir, img)
    iw, ih = Image.open(path).size
    ratio = iw / ih
    if ratio >= 2.2:
        # 结果条高度（chip + 单行 bullets）
        res_h = len(results) * 0.27 + 0.10
        res_y = 6.70 - res_h
        _place_image(s, path, 0.45, 2.24, 12.45, res_y - 0.12 - 2.24, caption)
        # 左侧 navy chip + 右侧 bullets
        add_rect(s, 0.45, res_y + 0.03, 1.5, res_h - 0.12, fill=NAVY,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.14)
        add_text(s, 0.45, res_y + 0.03, 1.5, res_h - 0.12,
                 [dict(runs=[("结果分析", "n")], size=12, color=WHITE, bold=True)],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        paras = [_result_bullet(b, size=11.5) for b in results]
        add_text(s, 2.2, res_y + 0.05, 10.7, res_h - 0.1, paras)
    else:
        _place_image(s, path, 0.45, 2.24, 8.1, 6.68 - 2.24, caption)
        paras = [_results_head(12.5)] + [_result_bullet(b, size=11.5) for b in results]
        add_text(s, 8.85, 2.34, 4.05, 4.2, paras)
    takeaway_bar(s, takeaway)
    page_number(s, idx)
    return s

def divider_slide(prs, idx, label, title_cn, title_en, question, agenda,
                  active=(0,), chain=None):
    """Results 小标题章节分隔页：纯文字排版(不放结果图，避免听众误解)。
    chain: 可选，证据链环节名列表(如 ["① 磁性来源", ...])，传入则在底部
    绘制进度条并用 active 索引高亮当前环节。"""
    s = new_slide(prs, bg=NAVY_DEEP)
    add_rect(s, 0.9, 1.55, 1.3, 0.05, fill=RED)
    add_text(s, 0.9, 0.8, 11.5, 0.6,
             [dict(runs=[(label, "n")], size=21, color=RED, bold=True)])
    add_text(s, 0.9, 1.85, 11.6, 0.85,
             [dict(runs=[(title_cn, "n")], size=33, color=WHITE, bold=True, line=1.15)])
    add_text(s, 0.9, 2.72, 11.6, 0.5,
             [dict(runs=[(title_en, "n")], size=14.5, color=SOFT_WHT, italic=True)])
    add_text(s, 0.9, 3.85, 11.5, 0.7,
             [dict(runs=[("本节问题   ", "n", 15), (question, "n", 19)],
                   size=19, color=WHITE, bold=True, line=1.2)])
    add_text(s, 0.9, 4.8, 11.5, 0.8,
             [dict(runs=[(agenda, "n")], size=14, color=SOFT_WHT, line=1.4)])
    # 底部：证据链进度（可选）
    if not chain:
        page_number(s, idx, dark_bg=True)
        return s
    add_rect(s, 0.9, 5.9, 11.53, 0.02, fill=RGBColor(0x3A, 0x4E, 0x70))
    n = len(chain)
    step_w = (11.53 - 0.21 * (n - 1)) / n
    for i, stage in enumerate(chain):
        cx = 0.9 + i * (step_w + 0.21)
        on = i in active
        add_rect(s, cx, 6.15, step_w, 0.62,
                 fill=(RED if on else None),
                 line=(None if on else RGBColor(0x3A, 0x4E, 0x70)), line_w=1.2,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.22)
        add_text(s, cx, 6.15, step_w, 0.62,
                 [dict(runs=[(stage, "n")], size=13,
                       color=(WHITE if on else SOFT_WHT), bold=on)],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    page_number(s, idx, dark_bg=True)
    return s
